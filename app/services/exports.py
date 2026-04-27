from __future__ import annotations

import base64
import mimetypes
from datetime import datetime
from pathlib import Path
from datetime import date

from jinja2 import Environment, FileSystemLoader, select_autoescape

from app.config import AppConfig
from app.models import ExportFormat, ExportRequest, ExportResult
from app.services.storage import StorageService
from markdown import markdown


class ExportService:
    def __init__(self, config: AppConfig, storage: StorageService) -> None:
        self.config = config
        self.storage = storage
        self.environment = Environment(
            loader=FileSystemLoader(str(config.templates_dir)),
            autoescape=select_autoescape(["html", "xml"]),
        )
        self.environment.filters["markdownify"] = lambda value: markdown(
            value or "",
            extensions=["fenced_code", "tables", "sane_lists"],
        )

    def run(self, request: ExportRequest) -> tuple[Path, list[ExportResult]]:
        stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        batch_dir = self.config.exports_dir / stamp
        batch_dir.mkdir(parents=True, exist_ok=True)

        html_paths: dict[str, Path] = {}
        results: list[ExportResult] = []
        for slug in request.project_slugs:
            loaded = self.storage.load_project(slug)
            if ExportFormat.HTML in request.formats or ExportFormat.PNG in request.formats or ExportFormat.PPTX in request.formats:
                html_path = self._write_html_export(batch_dir, slug, loaded)
                html_paths[slug] = html_path
                if ExportFormat.HTML in request.formats:
                    results.append(ExportResult(format=ExportFormat.HTML, output_path=str(html_path), success=True))
            if ExportFormat.PNG in request.formats:
                png_result = self._write_png_export(html_paths[slug], batch_dir / f"{slug}.png")
                results.append(png_result)

        if ExportFormat.PPTX in request.formats:
            pptx_result = self._write_pptx_export(
                batch_dir / "projstatus-summary.pptx",
                request.project_slugs,
                html_paths,
                batch_dir,
            )
            results.append(pptx_result)
        return batch_dir, results

    def _write_html_export(self, batch_dir: Path, slug: str, loaded) -> Path:
        template = self.environment.get_template("export_project.html")
        html = template.render(
            project=loaded.project,
            logo_data_uri=self._logo_data_uri(loaded.project),
            sections=loaded.sections,
            timeline_text=loaded.timeline_text,
            blocked_tasks=[task for task in loaded.project.tasks if task.column == "Blocked"],
            upcoming_milestones=sorted(
                [item for item in loaded.project.milestones if item.target_date],
                key=lambda item: item.target_date or date.max,
            ),
            export_mode=True,
        )
        output_path = batch_dir / f"{slug}.html"
        output_path.write_text(html, encoding="utf-8", newline="\n")
        return output_path

    def _write_png_export(self, html_path: Path, output_path: Path) -> ExportResult:
        try:
            from playwright.sync_api import sync_playwright
        except ImportError:
            return ExportResult(
                format=ExportFormat.PNG,
                output_path=str(output_path),
                success=False,
                message="Install the optional 'exports' dependency and run 'playwright install chromium' for PNG export.",
            )

        try:
            with sync_playwright() as playwright:
                browser = playwright.chromium.launch()
                page = browser.new_page(viewport={"width": 1440, "height": 1800})
                page.goto(html_path.as_uri(), wait_until="networkidle")
                page.screenshot(path=str(output_path), full_page=True)
                browser.close()
            return ExportResult(format=ExportFormat.PNG, output_path=str(output_path), success=True)
        except Exception as exc:  # pragma: no cover - depends on local browser install
            return ExportResult(format=ExportFormat.PNG, output_path=str(output_path), success=False, message=str(exc))

    def _write_pptx_export(
        self,
        output_path: Path,
        project_slugs: list[str],
        html_paths: dict[str, Path],
        batch_dir: Path,
    ) -> ExportResult:
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            return ExportResult(
                format=ExportFormat.PPTX,
                output_path=str(output_path),
                success=False,
                message="Install python-pptx to enable PowerPoint exports.",
            )

        try:
            presentation = Presentation()
            title_layout = presentation.slide_layouts[0]
            slide = presentation.slides.add_slide(title_layout)
            slide.shapes.title.text = "ProjStatus Summary"
            slide.placeholders[1].text = "Generated from local project files"

            for slug in project_slugs:
                loaded = self.storage.load_project(slug)
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                slide.shapes.title.text = loaded.project.name
                body = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(4.6), Inches(4.5)).text_frame
                body.text = f"Health: {loaded.project.health.value}"
                for line in [
                    f"Status: {loaded.project.status.value}",
                    f"Dates: {loaded.project.start_date or '-'} to {loaded.project.end_date or '-'}",
                    f"People: {len(loaded.project.people)}",
                    f"Milestones: {len(loaded.project.milestones)}",
                    f"Blocked tasks: {sum(1 for task in loaded.project.tasks if task.column == 'Blocked')}",
                ]:
                    body.add_paragraph().text = line

                hyperlink_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(5.0), Inches(0.6))
                run = hyperlink_box.text_frame.paragraphs[0].add_run()
                run.text = "Open interactive HTML export"
                run.hyperlink.address = html_paths[slug].as_uri()

                png_path = batch_dir / f"{slug}.png"
                if png_path.exists():
                    slide.shapes.add_picture(str(png_path), Inches(5.2), Inches(1.1), Inches(4.1))

            presentation.save(str(output_path))
            return ExportResult(format=ExportFormat.PPTX, output_path=str(output_path), success=True)
        except Exception as exc:
            return ExportResult(format=ExportFormat.PPTX, output_path=str(output_path), success=False, message=str(exc))

    def _logo_data_uri(self, project) -> str:
        logo_file = self.storage.resolve_logo_file(project)
        if logo_file is None:
            return ""
        mime_type = mimetypes.guess_type(logo_file.name)[0] or "application/octet-stream"
        encoded = base64.b64encode(logo_file.read_bytes()).decode("ascii")
        return f"data:{mime_type};base64,{encoded}"
